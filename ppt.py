

import os
import argparse
import fitz
import numpy as np
from PIL import Image
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
import openai
import time
import shutil
import warnings
import layoutparser as lp

# Suppress warnings
warnings.filterwarnings("ignore")

# Handle PIL compatibility
if not hasattr(Image, "LINEAR"):
    Image.LINEAR = Image.BILINEAR

# Required pip packages:
# pip install python-pptx PyMuPDF pdf2image pillow openai layoutparser detectron2 torch

def extract_text_from_pdf(pdf_path):
    """Extract text from each page of the PDF"""
    doc = fitz.open(pdf_path)
    texts = []
    for page in doc:
        texts.append(page.get_text())
    return texts


def generate_with_gpt(text, prompt_type="slide", page_number=1, is_first_page=False):
    """Generate content using GPT API"""
    # Configure your OpenAI API key
    openai.api_key = os.environ.get("OPENAI_API_KEY")
    
    if not openai.api_key:
        raise ValueError("Please set the OPENAI_API_KEY environment variable")
    
    # Define prompts for different content types
    prompts = {
        "slide": "Create 3-5 concise bullet points for a PowerPoint slide from the following text. "
                 "Each bullet point must be 10 words or less. Focus only on key concepts. "
                 "Format as bullet points with • symbol. "
                 "Also include a short, catchy title (up to 5 words) for this slide at the very beginning, "
                 "prefixed with 'TITLE: ':\n\n",
                 
        "speech": "Create concise teacher notes that correspond to slide #" + str(page_number) + " of a presentation. "
                  "These notes should help a teacher explain the slide's content and should include: "
                  "1) Key points to emphasize 2) Possible student questions 3) Clear explanations of complex concepts. "
                  "Keep it within 150-200 words and focus on clarity and teaching value:\n\n"
    }
    
    # Special handling for first page
    if is_first_page and prompt_type == "slide":
        prompts["slide"] = "This is the FIRST SLIDE of the presentation. Create a title slide with: " + \
                           "1) A main title (5-7 words maximum) prefixed with 'MAIN_TITLE: ' " + \
                           "2) A subtitle (10-15 words) prefixed with 'SUBTITLE: ' " + \
                           "Both should capture the essence of the document. Based on this text:\n\n"
    
    # Limit input text length to avoid token limits
    max_input_length = 3000
    text = text[:max_input_length]
    
    try:
        # Send request to OpenAI API
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",  # You can also use "gpt-4" for better quality
            messages=[
                {"role": "system", "content": "You are a helpful assistant that creates PowerPoint content."},
                {"role": "user", "content": prompts[prompt_type] + text}
            ],
            max_tokens=800,
            temperature=0.7
        )
        
        # Extract the generated content
        generated_content = response.choices[0].message.content.strip()
        return generated_content
    
    except Exception as e:
        print(f"Error generating content with GPT: {e}")
        # Fallback to basic extraction if GPT API fails
        if prompt_type == "slide":
            if is_first_page:
                return "MAIN_TITLE: Document Overview\nSUBTITLE: Key concepts and information"
            else:
                return "TITLE: Page " + str(page_number) + " Content\n" + extract_basic_bullet_points(text)
        else:
            return "Teacher notes for slide " + str(page_number) + ":\n\n" + text[:300] + "...\n\n[Note: Content truncated due to API error]"


def extract_basic_bullet_points(text, max_points=5):
    """Basic fallback function to extract bullet points without GPT"""
    sentences = text.split('. ')
    
    # Get first few sentences and format as bullet points
    points = sentences[:max_points]
    formatted_points = []
    
    for point in points:
        point = point.strip()
        if point:
            # Truncate long points to make them concise
            words = point.split()
            if len(words) > 8:
                point = ' '.join(words[:8]) + "..."
            formatted_points.append(f"• {point}")
    
    if not formatted_points:
        formatted_points.append("• No content available")
    
    return "\n".join(formatted_points)


def clean_title(title_text):
    """Remove asterisks from titles and clean up formatting"""
    return title_text.replace('*', '').strip()


def create_pptx(output_dir, pptx_filename="output_presentation.pptx", title_font_size=32, content_font_size=14):
    """Create PowerPoint presentation from processed content with adaptive layout"""
    prs = Presentation()
    
    # Change slide dimensions to 16:9 aspect ratio
    prs.slide_width = int(Inches(16))  # Set slide width to 16 inches
    prs.slide_height = int(Inches(9))  # Set slide height to 9 inches
    
    blank_layout = prs.slide_layouts[6]  # Blank layout
    title_layout = prs.slide_layouts[0]  # Title slide layout
    
    # Calculate slide dimensions and margins - IMPROVED CENTERING
    SLIDE_WIDTH = prs.slide_width / 914400
    SLIDE_HEIGHT = prs.slide_height / 914400
    left_margin = 1.0  # Increased for better centering
    right_margin = 1.0  # Increased for better centering
    top_margin = 1.5  # Increased spacing at the top
    bottom_margin = 0.8  # Slightly increased for 16:9
    title_spacing = 0.7  # Space between title and content
    gutter_in = 0.5  # Increased for better spacing
    usable_width = SLIDE_WIDTH - left_margin - right_margin
    usable_height = SLIDE_HEIGHT - top_margin - bottom_margin
    
    # Get and sort page folders
    page_folders = []
    for item in os.listdir(output_dir):
        if item.startswith("page_") and os.path.isdir(os.path.join(output_dir, item)):
            try:
                page_num = int(item.split("_")[1])
                page_folders.append((page_num, item))
            except:
                continue
    
    page_folders.sort()
    
    # Process each page folder
    for page_num, folder in page_folders:
        folder_path = os.path.join(output_dir, folder)
        is_first_page = (page_num == 1)
        
        # Read generated slide content
        slide_content_file = os.path.join(folder_path, "slide_script.txt")
        slide_content = ""
        slide_title = f"Page {page_num}"
        
        if os.path.exists(slide_content_file):
            try:
                with open(slide_content_file, "r", encoding="utf-8") as f:
                    slide_content = f.read()
                    
                    # Extract title from slide content
                    if is_first_page and "MAIN_TITLE:" in slide_content:
                        # Extract main title and subtitle for first slide
                        main_title_match = slide_content.split("MAIN_TITLE:", 1)[1].split("\n", 1)[0].strip()
                        main_title_match = clean_title(main_title_match)  # Clean title
                        
                        subtitle_match = ""
                        if "SUBTITLE:" in slide_content:
                            subtitle_match = slide_content.split("SUBTITLE:", 1)[1].split("\n", 1)[0].strip()
                            subtitle_match = clean_title(subtitle_match)  # Clean subtitle
                        
                        # Create title slide
                        title_slide = prs.slides.add_slide(title_layout)
                        title = title_slide.shapes.title
                        subtitle = title_slide.placeholders[1]
                        
                        # Add a background shape to maintain consistent visual style
                        background = title_slide.shapes.add_shape(
                            1,  # Rectangle
                            0, 0, prs.slide_width, int(Inches(1.2))
                        )
                        background.fill.solid()
                        background.fill.fore_color.rgb = RGBColor(0, 114, 198)  # Blue background for title
                        background.line.fill.background()  # No border
                        
                        # Send title shape to front
                        title.element.getparent().remove(title.element)
                        title_slide.shapes._spTree.append(title.element)
                        
                        # Send subtitle shape to front
                        subtitle.element.getparent().remove(subtitle.element)
                        title_slide.shapes._spTree.append(subtitle.element)
                        
                        # Position the title and subtitle better
                        title.top = int(Inches(2.5))  # Position title lower on the slide
                        title.left = int(Inches(1.0))  # Left margin for title
                        title.width = int(Inches(14.0))  # Width for title
                        title.height = int(Inches(1.5))  # Explicitly set height to ensure visibility
                        
                        subtitle.top = int(Inches(4.0))  # Position subtitle further down
                        subtitle.left = int(Inches(1.0))  # Left margin for subtitle  
                        subtitle.width = int(Inches(14.0))  # Width for subtitle
                        subtitle.height = int(Inches(1.2))  # Explicitly set height to ensure visibility
                        
                        # Set text first, then apply formatting for main title
                        title.text = main_title_match
                        title.text_frame.word_wrap = True
                        title.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  # Changed to ensure text is visible
                        
                        # Set text first, then apply formatting for subtitle
                        subtitle.text = subtitle_match
                        subtitle.text_frame.word_wrap = True
                        subtitle.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  # Changed to ensure text is visible
                        
                        # Apply custom font sizes
                        title.text_frame.paragraphs[0].runs[0].font.size = Pt(title_font_size + 8)  # Larger for better visibility
                        subtitle.text_frame.paragraphs[0].runs[0].font.size = Pt(content_font_size + 4)  # Larger for better visibility
                        
                        # Set text color to BLACK for better visibility
                        title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
                        subtitle.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
                        
                        # Skip further processing for title slide
                        continue
                    else:
                        # Extract slide title for regular slides
                        if "TITLE:" in slide_content:
                            slide_title = slide_content.split("TITLE:", 1)[1].split("\n", 1)[0].strip()
                            slide_title = clean_title(slide_title)  # Clean title
                            
                            # Remove the title line from content
                            slide_content = slide_content.split("\n", 1)[1] if "\n" in slide_content else ""
            except:
                slide_content = f"Content for page {page_num}"
        
        # Check for images
        image_files = []
        try:
            image_files = sorted([
                f for f in os.listdir(folder_path) 
                if f.startswith("image_") and f.endswith(".png")
            ])
        except:
            pass
        
        # Create content slide
        try:
            slide = prs.slides.add_slide(blank_layout)
            
            # Add title banner for better visual hierarchy
            title_banner = slide.shapes.add_shape(
                1,  # Rectangle
                0, 0, prs.slide_width, int(Inches(1.0))
            )
            title_banner.fill.solid()
            title_banner.fill.fore_color.rgb = RGBColor(0, 114, 198)  # Blue background for title
            title_banner.line.fill.background()  # No border
            
            # Add title with word wrap enabled - better centering
            title_box = slide.shapes.add_textbox(
                Inches(left_margin),
                Inches(0.25),  # Centered within the banner
                Inches(usable_width),
                Inches(0.7)
            )
            title_tf = title_box.text_frame
            title_tf.word_wrap = True
            title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # Automatically adjust text to fit shape
            title_tf.text = slide_title
            title_tf.paragraphs[0].alignment = 1  # Center alignment (0=left, 1=center, 2=right)
            title_tf.paragraphs[0].runs[0].font.size = Pt(title_font_size)
            title_tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text on blue banner
            
            # If no images, use full width for content with larger font - improved centering
            if not image_files:
                content_box = slide.shapes.add_textbox(
                    Inches(left_margin + 0.5),  # More indent for better centering
                    Inches(top_margin),  # Content starts after the title and spacing
                    Inches(usable_width - 1.0),  # Reduced width for better centering
                    Inches(usable_height)
                )
                content_tf = content_box.text_frame
                content_tf.text = slide_content
                content_tf.word_wrap = True
                content_tf.auto_size = MSO_AUTO_SIZE.NONE
                content_tf.vertical_anchor = MSO_ANCHOR.TOP
                
                # Set larger font size for slides with no images
                larger_content_font_size = content_font_size + 6  # Increase font size for better visibility on 16:9
                for paragraph in content_tf.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(larger_content_font_size)
            
            # If images available, add to slide with adaptive layout
            elif image_files:
                img_path = os.path.join(folder_path, image_files[0])
                
                try:
                    with Image.open(img_path) as pil_img:
                        orig_w, orig_h = pil_img.size
                    
                    # Calculate aspect ratio to determine layout
                    aspect_ratio = orig_h / orig_w
                    
                    # Wide image (flat) - use vertical layout (content on top, image below)
                    if aspect_ratio < 0.75:  # Width more than 1.33x height
                        # Content area - top section, better centered
                        content_height = usable_height * 0.40  # Reduced slightly for better spacing
                        content_box = slide.shapes.add_textbox(
                            Inches(left_margin + 0.5),  # More indent for better centering
                            Inches(top_margin),
                            Inches(usable_width - 1.0),  # Reduced width for better centering
                            Inches(content_height)
                        )
                        
                        # Image area - bottom section centered horizontally
                        pic_width_in = usable_width * 0.85  # Use 85% of width for better appearance
                        pic_height_in = pic_width_in * aspect_ratio
                        
                        # Ensure image isn't too tall
                        if pic_height_in > (usable_height - content_height - 0.3):  # Increased spacing
                            pic_height_in = usable_height - content_height - 0.3
                            pic_width_in = pic_height_in / aspect_ratio
                        
                        # Center the image horizontally
                        pic_left = left_margin + (usable_width - pic_width_in) / 2
                        pic_top = top_margin + content_height + 0.3  # Increased gap between content and image
                    
                    # Tall or square image - use horizontal layout (content left, image right)
                    else:
                        # For horizontal layout - adjusted for 16:9 with better centering
                        left_col_width = (usable_width - gutter_in) * 0.42  # Less width for text on 16:9
                        right_col_width = (usable_width - gutter_in) * 0.58  # More width for image on 16:9
                        
                        # Content area - left section with better centering
                        content_box = slide.shapes.add_textbox(
                            Inches(left_margin + 0.3),  # Slight indent for better centering
                            Inches(top_margin),
                            Inches(left_col_width - 0.3),  # Adjusted for indent
                            Inches(usable_height)
                        )
                        
                        # Image area - right section centered vertically
                        pic_width_in = right_col_width * 0.95  # Use 95% of column width
                        pic_height_in = pic_width_in * aspect_ratio
                        
                        # Ensure image isn't too tall
                        if pic_height_in > usable_height:
                            pic_height_in = usable_height * 0.95  # Use 95% of usable height
                            pic_width_in = pic_height_in / aspect_ratio
                        
                        pic_left = left_margin + left_col_width + gutter_in
                        # Center the image vertically
                        pic_top = top_margin + (usable_height - pic_height_in) / 2
                    
                    # Add content text
                    content_tf = content_box.text_frame
                    content_tf.text = slide_content
                    content_tf.word_wrap = True
                    content_tf.auto_size = MSO_AUTO_SIZE.NONE
                    content_tf.vertical_anchor = MSO_ANCHOR.TOP
                    
                    # Set font size - larger for 16:9
                    for paragraph in content_tf.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(content_font_size + 2)  # Slightly larger for 16:9
                    
                    # Add the image
                    slide.shapes.add_picture(
                        img_path,
                        Inches(pic_left),
                        Inches(pic_top),
                        width=Inches(pic_width_in),
                        height=Inches(pic_height_in)
                    )
                    
                except Exception as e:
                    print(f"Error processing image: {e}")
                    # Fallback to full-width content if image processing fails
                    content_box = slide.shapes.add_textbox(
                        Inches(left_margin + 0.5),  # More indent for better centering
                        Inches(top_margin),
                        Inches(usable_width - 1.0),  # Reduced width for better centering
                        Inches(usable_height)
                    )
                    content_tf = content_box.text_frame
                    content_tf.text = slide_content
                    content_tf.word_wrap = True
                    
                    # Set larger font size for fallback case
                    for paragraph in content_tf.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(content_font_size + 6)  # Increased font size for 16:9
        
        except Exception as e:
            print(f"Error creating slide for page {page_num}: {e}")
            # Add basic slide if complex layout fails
            try:
                slide = prs.slides.add_slide(blank_layout)
                text_box = slide.shapes.add_textbox(
                    Inches(left_margin + 0.5),  # More indent for better centering
                    Inches(top_margin),
                    Inches(usable_width - 1.0),  # Reduced width for better centering
                    Inches(usable_height)
                )
                text_tf = text_box.text_frame
                text_tf.text = f"Page {page_num}\n\n" + slide_content
                
                # Apply basic font size
                for paragraph in text_tf.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(content_font_size + 4)  # Increased for 16:9
            except:
                print(f"Failed to create even basic slide for page {page_num}")
    
    # Save presentation
    try:
        # Ensure output directory exists
        output_dir = os.path.dirname(pptx_filename)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir)
            
        prs.save(pptx_filename)
        print(f"PowerPoint presentation saved as '{pptx_filename}'.")
    except Exception as e:
        print(f"Error saving presentation: {e}")


def boxes_are_close(box1, box2, margin=10):
    """Check if two bounding boxes are close to each other"""
    x0, y0, x1, y1 = box1
    a0, b0, a1, b1 = box2
    gap_x = 0
    if x1 < a0:
        gap_x = a0 - x1
    elif a1 < x0:
        gap_x = x0 - a1
    gap_y = 0
    if y1 < b0:
        gap_y = b0 - y1
    elif b1 < y0:
        gap_y = y0 - b1
    return (gap_x <= margin) and (gap_y <= margin)


def union_box(box1, box2):
    """Merge two bounding boxes into one"""
    x0, y0, x1, y1 = box1
    a0, b0, a1, b1 = box2
    return (min(x0, a0), min(y0, b0), max(x1, a1), max(y1, b1))


def merge_rectangles(rectangles, margin=10):
    """Merge overlapping or nearby rectangles"""
    if not rectangles:
        return []
    merged = rectangles[:]
    changed = True
    while changed:
        changed = False
        new_merged = []
        used = [False] * len(merged)
        for i in range(len(merged)):
            if used[i]:
                continue
            current = merged[i]
            for j in range(i + 1, len(merged)):
                if used[j]:
                    continue
                if boxes_are_close(current, merged[j], margin):
                    current = union_box(current, merged[j])
                    used[j] = True
                    changed = True
            new_merged.append(current)
        merged = new_merged
    return merged


def process_pdf(pdf_path, output_dir):
    """Main function to process PDF and generate presentation content"""
    os.makedirs(output_dir, exist_ok=True)
    
    print("Converting PDF pages to images...")
    pages = convert_from_path(pdf_path, dpi=300)
    print(f"Converted {len(pages)} pages.")
    
    print("Extracting text from PDF...")
    texts = extract_text_from_pdf(pdf_path)
    
    # Load layout model directly from layoutparser as in reference code
    from layoutparser.models import Detectron2LayoutModel
    print("Loading layout detection model...")
    lp_model = Detectron2LayoutModel(
        "lp://PubLayNet/faster_rcnn_R_50_FPN_3x/config",
        extra_config=["MODEL.ROI_HEADS.SCORE_THRESH_TEST", 0.5],
        label_map={0: "Text", 1: "Title", 2: "List", 3: "Table", 4: "Figure"}
    )
    
    for i, page_image in enumerate(pages):
        page_number = i + 1
        is_first_page = (page_number == 1)
        
        print(f"\nProcessing page {page_number}...")
        page_folder = os.path.join(output_dir, f"page_{page_number}")
        os.makedirs(page_folder, exist_ok=True)
        
        # Get text for current page
        text = texts[i] if i < len(texts) else ""
        
        # Save original text
        with open(os.path.join(page_folder, "page_text.txt"), "w", encoding="utf-8") as f:
            f.write(text)
        
        # Generate slide content using GPT with page_number and first_page flag
        print(f"  Generating slide script for page {page_number}...")
        slide_content = generate_with_gpt(text, prompt_type="slide", page_number=page_number, is_first_page=is_first_page)
        
        # Generate teacher speech using GPT with page_number
        print(f"  Generating teacher script for page {page_number}...")
        teacher_script = generate_with_gpt(text, prompt_type="speech", page_number=page_number, is_first_page=is_first_page)
        
        # Save generated content
        with open(os.path.join(page_folder, "slide_script.txt"), "w", encoding="utf-8") as f:
            f.write(slide_content)
        
        with open(os.path.join(page_folder, "teacher_script.txt"), "w", encoding="utf-8") as f:
            f.write(teacher_script)
        
        # Extract images using exact image extraction code from reference
        try:
            # Convert PIL image to numpy array
            image_np = np.array(page_image)
            
            # Run layout detection
            layout = lp_model.detect(image_np)
            
            # Filter for figures and tables
            blocks = [b for b in layout if b.type in ["Figure", "Table"]]
            
            if blocks:
                # Extract coordinates and create rectangles
                rects = [tuple(map(int, b.coordinates)) for b in blocks]
                
                # Merge nearby rectangles
                merged_rects = merge_rectangles(rects, margin=20)
                
                # Save each detected region
                for j, rect in enumerate(merged_rects, start=1):
                    x0, y0, x1, y1 = rect
                    cropped = page_image.crop((x0, y0, x1, y1))
                    cropped.save(os.path.join(page_folder, f"image_{j}.png"))
                    print(f"    Saved image {j} from page {page_number}")
            else:
                print(f"    No figure/table blocks detected on page {page_number}.")
                # Don't save anything if no figures/tables detected
        except Exception as e:
            print(f"Error in layout detection: {e}")
        
        # Add small delay to avoid API rate limits
        time.sleep(1)
        
        print(f"Finished processing page {page_number}.")


def main():
    parser = argparse.ArgumentParser(description="Convert PDF to PowerPoint using GPT for content generation")
    parser.add_argument("pdf_path", help="Path to the PDF file")
    parser.add_argument("--output_dir", default="output", help="Output directory for processed files")
    parser.add_argument("--pptx", default="output_presentation.pptx", help="Output PowerPoint filename")
    parser.add_argument("--title_font_size", type=int, default=32, help="Font size for slide titles")
    parser.add_argument("--content_font_size", type=int, default=14, help="Font size for slide content")
    parser.add_argument("--widescreen", type=bool, default=True, help="Use 16:9 widescreen format (True) or 4:3 standard format (False)")
    args = parser.parse_args()
    
    if os.path.exists(args.output_dir):
        print(f"Cleaning up existing output directory: {args.output_dir}")
        shutil.rmtree(args.output_dir)
    
    process_pdf(args.pdf_path, args.output_dir)
    create_pptx(args.output_dir, pptx_filename=args.pptx, 
                title_font_size=args.title_font_size, 
                content_font_size=args.content_font_size)
    print("Processing complete.")


if __name__ == "__main__":
    main()