#!/usr/bin/env python3
"""
Extract each slide from a PPTX file as PNG images
Using multiple different methods to try extraction
"""

import os
import argparse
import subprocess
import tempfile
import shutil
from pathlib import Path
import sys
import platform

def extract_with_libreoffice(pptx_path, output_dir):
    """Extract slides using LibreOffice (for macOS and Linux)"""
    print("Trying to extract slides using LibreOffice...")
    
    try:
        # Check if LibreOffice is installed
        soffice_path = None
        if platform.system() == "Darwin":  # macOS
            possible_paths = [
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                "/Applications/OpenOffice.app/Contents/MacOS/soffice"
            ]
            for path in possible_paths:
                if os.path.exists(path):
                    soffice_path = path
                    break
        else:  # Linux or other systems
            try:
                soffice_path = subprocess.check_output(["which", "soffice"]).decode().strip()
            except:
                pass
            
        if not soffice_path:
            print("LibreOffice/OpenOffice not found")
            return False
            
        # Create temporary directory for PDF
        with tempfile.TemporaryDirectory() as temp_dir:
            # Convert PPTX to PDF
            cmd = [
                soffice_path,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", temp_dir,
                str(pptx_path)
            ]
            
            subprocess.run(cmd, check=True, capture_output=True)
            
            # Find the generated PDF file
            pdf_file = None
            for file in os.listdir(temp_dir):
                if file.endswith(".pdf"):
                    pdf_file = os.path.join(temp_dir, file)
                    break
                    
            if not pdf_file:
                print("Conversion failed: PDF file not generated")
                return False
                
            # Use pdftoppm to convert PDF to PNG images
            try:
                pdftoppm_cmd = ["pdftoppm", "-png", pdf_file, os.path.join(output_dir, "slide")]
                subprocess.run(pdftoppm_cmd, check=True)
                
                # Rename generated files
                for i, file in enumerate(sorted(os.listdir(output_dir))):
                    if file.startswith("slide-") and file.endswith(".png"):
                        old_path = os.path.join(output_dir, file)
                        new_path = os.path.join(output_dir, f"slide_{i+1}.png")
                        os.rename(old_path, new_path)
                
                return True
            except:
                print("pdftoppm failed, trying other methods...")
                return False
    except Exception as e:
        print(f"LibreOffice extraction failed: {e}")
        return False

def extract_with_pdf2image(pptx_path, output_dir):
    """Extract slides using pdf2image library (requires converting PPTX to PDF first)"""
    print("Trying to extract slides using pdf2image...")
    
    try:
        # First try to import necessary libraries
        try:
            from pdf2image import convert_from_path
        except ImportError:
            print("Please install pdf2image: pip install pdf2image")
            return False
        
        # Create temporary directory for PDF
        with tempfile.TemporaryDirectory() as temp_dir:
            # Try multiple methods to convert to PDF
            pdf_file = None
            
            # 1. Try using LibreOffice
            try:
                soffice_path = None
                if platform.system() == "Darwin":  # macOS
                    possible_paths = [
                        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                        "/Applications/OpenOffice.app/Contents/MacOS/soffice"
                    ]
                    for path in possible_paths:
                        if os.path.exists(path):
                            soffice_path = path
                            break
                else:  # Linux or other systems
                    try:
                        soffice_path = subprocess.check_output(["which", "soffice"]).decode().strip()
                    except:
                        pass
                        
                if soffice_path:
                    cmd = [
                        soffice_path,
                        "--headless",
                        "--convert-to", "pdf",
                        "--outdir", temp_dir,
                        str(pptx_path)
                    ]
                    
                    subprocess.run(cmd, check=True, capture_output=True)
                    
                    # Find the generated PDF file
                    for file in os.listdir(temp_dir):
                        if file.endswith(".pdf"):
                            pdf_file = os.path.join(temp_dir, file)
                            break
            except:
                pass
            
            # 2. If LibreOffice fails, try using unoconv
            if not pdf_file:
                try:
                    unoconv_path = subprocess.check_output(["which", "unoconv"]).decode().strip()
                    if unoconv_path:
                        cmd = [
                            unoconv_path,
                            "-f", "pdf",
                            "-o", temp_dir,
                            str(pptx_path)
                        ]
                        
                        subprocess.run(cmd, check=True, capture_output=True)
                        
                        # Find the generated PDF file
                        for file in os.listdir(temp_dir):
                            if file.endswith(".pdf"):
                                pdf_file = os.path.join(temp_dir, file)
                                break
                except:
                    pass
            
            if not pdf_file:
                print("Failed to convert to PDF")
                return False
                
            # Use pdf2image to extract slides
            pages = convert_from_path(pdf_file)
            
            # Save each page as PNG
            for i, page in enumerate(pages):
                page.save(os.path.join(output_dir, f"slide_{i+1}.png"), "PNG")
                
            return True
    except Exception as e:
        print(f"pdf2image extraction failed: {e}")
        return False

def extract_with_python_pptx_screenshots(pptx_path, output_dir):
    """Extract slide content using python-pptx and create better visualizations"""
    print("Trying to extract slides using python-pptx...")
    
    try:
        # Import necessary libraries
        from pptx import Presentation
        from PIL import Image, ImageDraw, ImageFont
        import io
        
        # Load presentation
        prs = Presentation(pptx_path)
        
        # For each slide, gather content and create visualization
        for i, slide in enumerate(prs.slides):
            # Create a white background image
            img = Image.new('RGB', (1280, 720), color='white')
            draw = ImageDraw.Draw(img)
            
            # Track positions for elements
            current_y = 50
            
            # Try to extract slide title
            title = ""
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    if shape.has_text_frame and shape.text_frame.text:
                        title = shape.text_frame.text
                        break
            
            # Try to extract slide elements and content
            content_items = []
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    text = shape.text_frame.text
                    if text and text != title:  # Don't repeat the title
                        content_items.append(text)
                
                # Try to extract images if available
                if hasattr(shape, "image") and shape.image:
                    try:
                        # For PPTX, image data is often accessible
                        img_stream = io.BytesIO(shape.image.blob)
                        slide_img = Image.open(img_stream)
                        
                        # Resize while maintaining aspect ratio
                        max_width = 800
                        max_height = 500
                        img_width, img_height = slide_img.size
                        
                        # Calculate new dimensions
                        ratio = min(max_width/img_width, max_height/img_height)
                        new_width = int(img_width * ratio)
                        new_height = int(img_height * ratio)
                        
                        # Resize and paste onto our image
                        slide_img = slide_img.resize((new_width, new_height), Image.LANCZOS)
                        img.paste(slide_img, (1280//2 - new_width//2, current_y + 200))
                        
                        # Update position
                        current_y += new_height + 20
                    except Exception as img_error:
                        print(f"Failed to extract image: {img_error}")
            
            # Get a decent font
            try:
                # Try to get a nice font, falling back to default if needed
                font_large = ImageFont.truetype("Arial", 40)
                font_medium = ImageFont.truetype("Arial", 30)
                font_small = ImageFont.truetype("Arial", 24)
            except:
                font_large = ImageFont.load_default()
                font_medium = font_large
                font_small = font_large
                
            # Draw slide number
            draw.text((1150, 30), f"Slide {i+1}", fill='black', font=font_medium)
            
            # Draw title centered
            if title:
                # Center the title text
                w, h = draw.textlength(title, font=font_large), font_large.size
                draw.text((640 - w/2, 50), title, fill='black', font=font_large)
                current_y = 120
            
            # Draw content
            for idx, content in enumerate(content_items):
                # Limit each content item to reasonable length
                content_lines = []
                words = content.split()
                current_line = ""
                
                # Handle multi-line content with wrapping
                for word in words:
                    if len(current_line + " " + word) * 10 < 1000:  # Approximate width check
                        current_line += (" " + word if current_line else word)
                    else:
                        content_lines.append(current_line)
                        current_line = word
                
                if current_line:
                    content_lines.append(current_line)
                
                # Draw each line
                for line in content_lines[:15]:  # Limit to 15 lines
                    draw.text((100, current_y), line, fill='black', font=font_small)
                    current_y += 30
                
                # Add spacing between content items
                current_y += 10
                
                # Don't let content go beyond the bottom
                if current_y > 680:
                    draw.text((100, 680), "... more content not shown ...", fill='red', font=font_small)
                    break
            
            # Save the image
            img.save(os.path.join(output_dir, f"slide_{i+1}.png"))
            
        return True
    except Exception as e:
        print(f"python-pptx extraction failed: {e}")
        return False

def extract_with_comtypes(pptx_path, output_dir):
    """Extract slides using comtypes to control PowerPoint (Windows only)"""
    if platform.system() != "Windows":
        print("comtypes method is Windows-only")
        return False
        
    print("Trying to extract slides using comtypes...")
    
    try:
        import comtypes.client
        
        # Get absolute path of PowerPoint file
        pptx_path = os.path.abspath(pptx_path)
        
        # Initialize PowerPoint application
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = True
        
        # Open presentation
        presentation = powerpoint.Presentations.Open(pptx_path)
        
        # Set output path
        output_path = os.path.join(os.path.abspath(output_dir), "slide")
        
        # Save slides as PNG images
        presentation.SaveAs(output_path, 18)  # 18 is the value for PNG format
        
        # Close presentation and PowerPoint
        presentation.Close()
        powerpoint.Quit()
        
        return True
    except Exception as e:
        print(f"comtypes extraction failed: {e}")
        return False

def main():
    parser = argparse.ArgumentParser(description="Extract slides from PPTX as PNG images")
    parser.add_argument("pptx_path", help="Path to PPTX file")
    parser.add_argument("--output_dir", default="slides", help="Output directory")
    
    args = parser.parse_args()
    
    # Ensure output directory exists
    os.makedirs(args.output_dir, exist_ok=True)
    
    # Try different extraction methods
    success = False
    
    # 1. If on Windows, try comtypes first
    if platform.system() == "Windows":
        success = extract_with_comtypes(args.pptx_path, args.output_dir)
    
    # 2. Try LibreOffice method
    if not success:
        success = extract_with_libreoffice(args.pptx_path, args.output_dir)
    
    # 3. Try pdf2image method
    if not success:
        success = extract_with_pdf2image(args.pptx_path, args.output_dir)
    
    # 4. Finally try python-pptx method
    if not success:
        success = extract_with_python_pptx_screenshots(args.pptx_path, args.output_dir)
    
    # Check if slides were extracted
    slides = [f for f in os.listdir(args.output_dir) if f.startswith("slide_") and f.endswith(".png")]
    
    if slides:
        print(f"Successfully extracted {len(slides)} slide images to {args.output_dir}")
    else:
        print("Failed to extract any slides. Try installing LibreOffice or run on Windows.")
        print("Possible solutions:")
        print("1. Install LibreOffice: https://www.libreoffice.org/download/")
        print("2. Install pdf2image: pip install pdf2image")
        print("3. On Windows, ensure Microsoft PowerPoint is installed")
        print("4. Manually open the PPTX and save as PNG images")

if __name__ == "__main__":
    main()