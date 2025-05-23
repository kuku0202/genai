# import os
# import argparse
# import fitz
# import numpy as np
# from PIL import Image
# from pdf2image import convert_from_path
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
# from pptx.dml.color import RGBColor
# import openai
# import time
# import shutil
# print("📌 pdftoppm path:", shutil.which("pdftoppm"))
# import warnings
# import layoutparser as lp

# # Suppress warnings
# warnings.filterwarnings("ignore")

# # Handle PIL compatibility
# if not hasattr(Image, "LINEAR"):
#     Image.LINEAR = Image.BILINEAR

# # Required pip packages:
# # pip install python-pptx teach pdf2image pillow openai layoutparser detectron2 torch

# def extract_text_from_pdf(pdf_path):
#     """Extract text from each page of the PDF"""
#     doc = fitz.open(pdf_path)
#     texts = []
#     for page in doc:
#         texts.append(page.get_text())
#     return texts


# def generate_with_gpt(text, prompt_type="slide", page_number=1, is_first_page=False):
#     """Generate content using GPT API"""
#     # Configure your OpenAI API key
#     openai.api_key = os.environ.get("OPENAI_API_KEY")
    
#     if not openai.api_key:
#         raise ValueError("Please set the OPENAI_API_KEY environment variable")
    
#     # Define prompts for different content types
#     prompts = {
#         "slide": "Create 3-5 concise bullet points for a PowerPoint slide from the following text. "
#                  "Each bullet point must be 10 words or less. Focus only on key concepts. "
#                  "Format as bullet points with • symbol. "
#                  "Also include a short, catchy title (up to 5 words) for this slide at the very beginning, "
#                  "prefixed with 'TITLE: ':\n\n",
                 
#         # "speech": "Write a clear, conversational script for slide #" + str(page_number) + " that will be read aloud verbatim. "
#         #   "Focus on explaining the slide's content in a natural speaking voice. "
#         #   "Include explanations of key concepts, highlight important points, and anticipate common questions. "
#         #   "Keep it within 150-200 words, use a direct, engaging tone, and avoid any meta-references or instructions. "
#         #   "The script should flow naturally when read aloud:\n\n"
#         "speech": "Write a conversational lecture script for this slide as if you're an experienced professor speaking to your class. Build naturally on concepts covered earlier in your lecture and maintain a logical progression of ideas. Use a warm, natural teaching voice with occasional rhetorical questions, thoughtful pauses, and authentic phrases ('Now, what's interesting here is...', 'Let me emphasize this point...', 'You might be wondering...'). Include brief anecdotes or examples where appropriate, connect to previously explained ideas, and anticipate student questions. Vary your sentence structure with both short, punchy statements and more complex explanations. Limit to 150-200 words, maintain a clear focus on the slide content, and ensure it flows naturally when read aloud. The script should sound like authentic classroom teaching that fits seamlessly into the overall lecture:\n\n"
#     }
    
#     # Special handling for first page
#     if is_first_page and prompt_type == "slide":
#         prompts["slide"] = "This is the FIRST SLIDE of the presentation. Create a title slide with: " + \
#                            "1) A main title (5-7 words maximum) prefixed with 'MAIN_TITLE: ' " + \
#                            "2) A subtitle (10-15 words) prefixed with 'SUBTITLE: ' " + \
#                            "Both should capture the essence of the document. Based on this text:\n\n"
    
#     # Limit input text length to avoid token limits
#     max_input_length = 3000
#     text = text[:max_input_length]
    
#     try:
#         # Send request to OpenAI API
#         response = openai.ChatCompletion.create(
#             model="gpt-3.5-turbo",  # You can also use "gpt-4" for better quality
#             messages=[
#                 {"role": "system", "content": "You are a helpful assistant that creates PowerPoint content."},
#                 {"role": "user", "content": prompts[prompt_type] + text}
#             ],
#             max_tokens=800,
#             temperature=0.7
#         )
        
#         # Extract the generated content
#         generated_content = response.choices[0].message.content.strip()
#         return generated_content
    
#     except Exception as e:
#         print(f"Error generating content with GPT: {e}")
#         # Fallback to basic extraction if GPT API fails
#         if prompt_type == "slide":
#             if is_first_page:
#                 return "MAIN_TITLE: Document Overview\nSUBTITLE: Key concepts and information"
#             else:
#                 return "TITLE: Page " + str(page_number) + " Content\n" + extract_basic_bullet_points(text)
#         else:
#             return "Teacher notes for slide " + str(page_number) + ":\n\n" + text[:300] + "...\n\n[Note: Content truncated due to API error]"


# def extract_basic_bullet_points(text, max_points=5):
#     """Basic fallback function to extract bullet points without GPT"""
#     sentences = text.split('. ')
    
#     # Get first few sentences and format as bullet points
#     points = sentences[:max_points]
#     formatted_points = []
    
#     for point in points:
#         point = point.strip()
#         if point:
#             # Truncate long points to make them concise
#             words = point.split()
#             if len(words) > 8:
#                 point = ' '.join(words[:8]) + "..."
#             formatted_points.append(f"• {point}")
    
#     if not formatted_points:
#         formatted_points.append("• No content available")
    
#     return "\n".join(formatted_points)


# def clean_title(title_text):
#     """Remove asterisks from titles and clean up formatting"""
#     return title_text.replace('*', '').strip()


# def create_pptx(output_dir, pptx_filename="output_presentation.pptx", title_font_size=32, content_font_size=14):
#     """Create PowerPoint presentation from processed content with adaptive layout"""
#     prs = Presentation()
    
#     # Change slide dimensions to 16:9 aspect ratio
#     prs.slide_width = int(Inches(16))  # Set slide width to 16 inches
#     prs.slide_height = int(Inches(9))  # Set slide height to 9 inches
    
#     blank_layout = prs.slide_layouts[6]  # Blank layout
#     title_layout = prs.slide_layouts[0]  # Title slide layout
    
#     # Calculate slide dimensions and margins - IMPROVED CENTERING
#     SLIDE_WIDTH = prs.slide_width / 914400
#     SLIDE_HEIGHT = prs.slide_height / 914400
#     left_margin = 1.0  # Increased for better centering
#     right_margin = 1.0  # Increased for better centering
#     top_margin = 1.5  # Increased spacing at the top
#     bottom_margin = 0.8  # Slightly increased for 16:9
#     title_spacing = 0.7  # Space between title and content
#     gutter_in = 0.5  # Increased for better spacing
#     usable_width = SLIDE_WIDTH - left_margin - right_margin
#     usable_height = SLIDE_HEIGHT - top_margin - bottom_margin
    
#     # Get and sort page folders
#     page_folders = []
#     for item in os.listdir(output_dir):
#         if item.startswith("page_") and os.path.isdir(os.path.join(output_dir, item)):
#             try:
#                 page_num = int(item.split("_")[1])
#                 page_folders.append((page_num, item))
#             except:
#                 continue
    
#     page_folders.sort()
    
#     # Process each page folder
#     for page_num, folder in page_folders:
#         folder_path = os.path.join(output_dir, folder)
#         is_first_page = (page_num == 1)
        
#         # Read generated slide content
#         slide_content_file = os.path.join(folder_path, "slide_script.txt")
#         slide_content = ""
#         slide_title = f"Page {page_num}"
        
#         if os.path.exists(slide_content_file):
#             try:
#                 with open(slide_content_file, "r", encoding="utf-8") as f:
#                     slide_content = f.read()
                    
#                     # Extract title from slide content
#                     if is_first_page and "MAIN_TITLE:" in slide_content:
#                         # Extract main title and subtitle for first slide
#                         main_title_match = slide_content.split("MAIN_TITLE:", 1)[1].split("\n", 1)[0].strip()
#                         main_title_match = clean_title(main_title_match)  # Clean title
                        
#                         subtitle_match = ""
#                         if "SUBTITLE:" in slide_content:
#                             subtitle_match = slide_content.split("SUBTITLE:", 1)[1].split("\n", 1)[0].strip()
#                             subtitle_match = clean_title(subtitle_match)  # Clean subtitle
                        
#                         # Create title slide
#                         title_slide = prs.slides.add_slide(title_layout)
#                         title = title_slide.shapes.title
#                         subtitle = title_slide.placeholders[1]
                        
#                         # Add a background shape to maintain consistent visual style
#                         background = title_slide.shapes.add_shape(
#                             1,  # Rectangle
#                             0, 0, prs.slide_width, int(Inches(1.2))
#                         )
#                         background.fill.solid()
#                         background.fill.fore_color.rgb = RGBColor(0, 114, 198)  # Blue background for title
#                         background.line.fill.background()  # No border
                        
#                         # Send title shape to front
#                         title.element.getparent().remove(title.element)
#                         title_slide.shapes._spTree.append(title.element)
                        
#                         # Send subtitle shape to front
#                         subtitle.element.getparent().remove(subtitle.element)
#                         title_slide.shapes._spTree.append(subtitle.element)
                        
#                         # Position the title and subtitle better
#                         title.top = int(Inches(2.5))  # Position title lower on the slide
#                         title.left = int(Inches(1.0))  # Left margin for title
#                         title.width = int(Inches(14.0))  # Width for title
#                         title.height = int(Inches(1.5))  # Explicitly set height to ensure visibility
                        
#                         subtitle.top = int(Inches(4.0))  # Position subtitle further down
#                         subtitle.left = int(Inches(1.0))  # Left margin for subtitle  
#                         subtitle.width = int(Inches(14.0))  # Width for subtitle
#                         subtitle.height = int(Inches(1.2))  # Explicitly set height to ensure visibility
                        
#                         # Set text first, then apply formatting for main title
#                         title.text = main_title_match
#                         title.text_frame.word_wrap = True
#                         title.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  # Changed to ensure text is visible
                        
#                         # Set text first, then apply formatting for subtitle
#                         subtitle.text = subtitle_match
#                         subtitle.text_frame.word_wrap = True
#                         subtitle.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  # Changed to ensure text is visible
                        
#                         # Apply custom font sizes
#                         title.text_frame.paragraphs[0].runs[0].font.size = Pt(title_font_size + 8)  # Larger for better visibility
#                         subtitle.text_frame.paragraphs[0].runs[0].font.size = Pt(content_font_size + 4)  # Larger for better visibility
                        
#                         # Set text color to BLACK for better visibility
#                         title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
#                         subtitle.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
                        
#                         # Skip further processing for title slide
#                         continue
#                     else:
#                         # Extract slide title for regular slides
#                         if "TITLE:" in slide_content:
#                             slide_title = slide_content.split("TITLE:", 1)[1].split("\n", 1)[0].strip()
#                             slide_title = clean_title(slide_title)  # Clean title
                            
#                             # Remove the title line from content
#                             slide_content = slide_content.split("\n", 1)[1] if "\n" in slide_content else ""
#             except:
#                 slide_content = f"Content for page {page_num}"
        
#         # Check for images
#         image_files = []
#         try:
#             image_files = sorted([
#                 f for f in os.listdir(folder_path) 
#                 if f.startswith("image_") and f.endswith(".png")
#             ])
#         except:
#             pass
        
#         # Create content slide
#         try:
#             slide = prs.slides.add_slide(blank_layout)
            
#             # Add title banner for better visual hierarchy
#             title_banner = slide.shapes.add_shape(
#                 1,  # Rectangle
#                 0, 0, prs.slide_width, int(Inches(1.0))
#             )
#             title_banner.fill.solid()
#             title_banner.fill.fore_color.rgb = RGBColor(0, 114, 198)  # Blue background for title
#             title_banner.line.fill.background()  # No border
            
#             # Add title with word wrap enabled - better centering
#             title_box = slide.shapes.add_textbox(
#                 Inches(left_margin),
#                 Inches(0.25),  # Centered within the banner
#                 Inches(usable_width),
#                 Inches(0.7)
#             )
#             title_tf = title_box.text_frame
#             title_tf.word_wrap = True
#             title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # Automatically adjust text to fit shape
#             title_tf.text = slide_title
#             title_tf.paragraphs[0].alignment = 1  # Center alignment (0=left, 1=center, 2=right)
#             title_tf.paragraphs[0].runs[0].font.size = Pt(title_font_size)
#             title_tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text on blue banner
            
#             # If no images, use full width for content with larger font - improved centering
#             if not image_files:
#                 content_box = slide.shapes.add_textbox(
#                     Inches(left_margin + 0.5),  # More indent for better centering
#                     Inches(top_margin),  # Content starts after the title and spacing
#                     Inches(usable_width - 1.0),  # Reduced width for better centering
#                     Inches(usable_height)
#                 )
#                 content_tf = content_box.text_frame
#                 content_tf.text = slide_content
#                 content_tf.word_wrap = True
#                 content_tf.auto_size = MSO_AUTO_SIZE.NONE
#                 content_tf.vertical_anchor = MSO_ANCHOR.TOP
                
#                 # Set larger font size for slides with no images
#                 larger_content_font_size = content_font_size + 6  # Increase font size for better visibility on 16:9
#                 for paragraph in content_tf.paragraphs:
#                     for run in paragraph.runs:
#                         run.font.size = Pt(larger_content_font_size)
            
#             # If images available, add to slide with adaptive layout
#             elif image_files:
#                 img_path = os.path.join(folder_path, image_files[0])
                
#                 try:
#                     with Image.open(img_path) as pil_img:
#                         orig_w, orig_h = pil_img.size
#                         # if orig_w < 300 or orig_h < 200:
#                         #     print("Image too small, skipping to avoid poor quality.")
#                         #     continue
#                     # Calculate aspect ratio to determine layout
#                     aspect_ratio = orig_h / orig_w
                    
#                     # Wide image (flat) - use vertical layout (content on top, image below)
#                     if aspect_ratio < 0.75:  # Width more than 1.33x height
#                         # Content area - top section, better centered
#                         content_height = usable_height * 0.40  # Reduced slightly for better spacing
#                         content_box = slide.shapes.add_textbox(
#                             Inches(left_margin + 0.5),  # More indent for better centering
#                             Inches(top_margin),
#                             Inches(usable_width - 1.0),  # Reduced width for better centering
#                             Inches(content_height)
#                         )
                        
#                         # Image area - bottom section centered horizontally
#                         pic_width_in = usable_width * 0.85  # Use 85% of width for better appearance
#                         pic_height_in = pic_width_in * aspect_ratio
                        
#                         # Ensure image isn't too tall
#                         if pic_height_in > (usable_height - content_height - 0.3):  # Increased spacing
#                             pic_height_in = usable_height - content_height - 0.3
#                             pic_width_in = pic_height_in / aspect_ratio
                        
#                         # Center the image horizontally
#                         pic_left = left_margin + (usable_width - pic_width_in) / 2
#                         pic_top = top_margin + content_height + 0.3  # Increased gap between content and image
                    
#                     # Tall or square image - use horizontal layout (content left, image right)
#                     else:
#                         # For horizontal layout - adjusted for 16:9 with better centering
#                         left_col_width = (usable_width - gutter_in) * 0.42  # Less width for text on 16:9
#                         right_col_width = (usable_width - gutter_in) * 0.58  # More width for image on 16:9
                        
#                         # Content area - left section with better centering
#                         content_box = slide.shapes.add_textbox(
#                             Inches(left_margin + 0.3),  # Slight indent for better centering
#                             Inches(top_margin),
#                             Inches(left_col_width - 0.3),  # Adjusted for indent
#                             Inches(usable_height)
#                         )
                        
#                         # Image area - right section centered vertically
#                         pic_width_in = right_col_width * 0.95  # Use 95% of column width
#                         pic_height_in = pic_width_in * aspect_ratio
                        
#                         # Ensure image isn't too tall
#                         if pic_height_in > usable_height:
#                             pic_height_in = usable_height * 0.95  # Use 95% of usable height
#                             pic_width_in = pic_height_in / aspect_ratio
                        
#                         pic_left = left_margin + left_col_width + gutter_in
#                         # Center the image vertically
#                         pic_top = top_margin + (usable_height - pic_height_in) / 2
                    
#                     # Add content text
#                     content_tf = content_box.text_frame
#                     content_tf.text = slide_content
#                     content_tf.word_wrap = True
#                     content_tf.auto_size = MSO_AUTO_SIZE.NONE
#                     content_tf.vertical_anchor = MSO_ANCHOR.TOP
                    
#                     # Set font size - larger for 16:9
#                     for paragraph in content_tf.paragraphs:
#                         for run in paragraph.runs:
#                             run.font.size = Pt(content_font_size + 2)  # Slightly larger for 16:9
                    
#                     # Add the image
#                     slide.shapes.add_picture(
#                         img_path,
#                         Inches(pic_left),
#                         Inches(pic_top),
#                         width=Inches(pic_width_in),
#                         height=Inches(pic_height_in)
#                     )
                    
#                 except Exception as e:
#                     print(f"Error processing image: {e}")
#                     # Fallback to full-width content if image processing fails
#                     content_box = slide.shapes.add_textbox(
#                         Inches(left_margin + 0.5),  # More indent for better centering
#                         Inches(top_margin),
#                         Inches(usable_width - 1.0),  # Reduced width for better centering
#                         Inches(usable_height)
#                     )
#                     content_tf = content_box.text_frame
#                     content_tf.text = slide_content
#                     content_tf.word_wrap = True
                    
#                     # Set larger font size for fallback case
#                     for paragraph in content_tf.paragraphs:
#                         for run in paragraph.runs:
#                             run.font.size = Pt(content_font_size + 6)  # Increased font size for 16:9
        
#         except Exception as e:
#             print(f"Error creating slide for page {page_num}: {e}")
#             # Add basic slide if complex layout fails
#             try:
#                 slide = prs.slides.add_slide(blank_layout)
#                 text_box = slide.shapes.add_textbox(
#                     Inches(left_margin + 0.5),  # More indent for better centering
#                     Inches(top_margin),
#                     Inches(usable_width - 1.0),  # Reduced width for better centering
#                     Inches(usable_height)
#                 )
#                 text_tf = text_box.text_frame
#                 text_tf.text = f"Page {page_num}\n\n" + slide_content
                
#                 # Apply basic font size
#                 for paragraph in text_tf.paragraphs:
#                     for run in paragraph.runs:
#                         run.font.size = Pt(content_font_size + 4)  # Increased for 16:9
#             except:
#                 print(f"Failed to create even basic slide for page {page_num}")
    
#     # Save presentation
#     try:
#         # Ensure output directory exists
#         output_dir = os.path.dirname(pptx_filename)
#         if output_dir and not os.path.exists(output_dir):
#             os.makedirs(output_dir)
            
#         prs.save(pptx_filename)
#         print(f"PowerPoint presentation saved as '{pptx_filename}'.")
#     except Exception as e:
#         print(f"Error saving presentation: {e}")


# def boxes_are_close(box1, box2, margin=10):
#     """Check if two bounding boxes are close to each other"""
#     x0, y0, x1, y1 = box1
#     a0, b0, a1, b1 = box2
#     gap_x = 0
#     if x1 < a0:
#         gap_x = a0 - x1
#     elif a1 < x0:
#         gap_x = x0 - a1
#     gap_y = 0
#     if y1 < b0:
#         gap_y = b0 - y1
#     elif b1 < y0:
#         gap_y = y0 - b1
#     return (gap_x <= margin) and (gap_y <= margin)


# def union_box(box1, box2):
#     """Merge two bounding boxes into one"""
#     x0, y0, x1, y1 = box1
#     a0, b0, a1, b1 = box2
#     return (min(x0, a0), min(y0, b0), max(x1, a1), max(y1, b1))


# def merge_rectangles(rectangles, margin=10):
#     """Merge overlapping or nearby rectangles"""
#     if not rectangles:
#         return []
#     merged = rectangles[:]
#     changed = True
#     while changed:
#         changed = False
#         new_merged = []
#         used = [False] * len(merged)
#         for i in range(len(merged)):
#             if used[i]:
#                 continue
#             current = merged[i]
#             for j in range(i + 1, len(merged)):
#                 if used[j]:
#                     continue
#                 if boxes_are_close(current, merged[j], margin):
#                     current = union_box(current, merged[j])
#                     used[j] = True
#                     changed = True
#             new_merged.append(current)
#         merged = new_merged
#     return merged


# def process_pdf(pdf_path, output_dir):
#     """Main function to process PDF and generate presentation content"""
#     os.makedirs(output_dir, exist_ok=True)
    
#     print("Converting PDF pages to images...")
#     pages = convert_from_path(pdf_path, dpi=300, poppler_path="/usr/bin")
#     print(f"Converted {len(pages)} pages.")
    
#     print("Extracting text from PDF...")
#     texts = extract_text_from_pdf(pdf_path)
    
#     # Load layout model directly from layoutparser as in reference code
#     from layoutparser.models import Detectron2LayoutModel
#     print("Loading layout detection model...")
#     lp_model = Detectron2LayoutModel(
#         "lp://PubLayNet/faster_rcnn_R_50_FPN_3x/config",
#         model_path="/models/layoutparser/publaynet/model_final.pth",
#         extra_config=["MODEL.ROI_HEADS.SCORE_THRESH_TEST", 0.5],
#         label_map={0: "Text", 1: "Title", 2: "List", 3: "Table", 4: "Figure"}
#     )
    
#     for i, page_image in enumerate(pages):
#         page_number = i + 1
#         is_first_page = (page_number == 1)
        
#         print(f"\nProcessing page {page_number}...")
#         page_folder = os.path.join(output_dir, f"page_{page_number}")
#         os.makedirs(page_folder, exist_ok=True)
        
#         # Get text for current page
#         text = texts[i] if i < len(texts) else ""
        
#         # Save original text
#         with open(os.path.join(page_folder, "page_text.txt"), "w", encoding="utf-8") as f:
#             f.write(text)
        
#         # Generate slide content using GPT with page_number and first_page flag
#         print(f"  Generating slide script for page {page_number}...")
#         slide_content = generate_with_gpt(text, prompt_type="slide", page_number=page_number, is_first_page=is_first_page)
        
#         # Generate teacher speech using GPT with page_number
#         print(f"  Generating teacher script for page {page_number}...")
#         teacher_script = generate_with_gpt(text, prompt_type="speech", page_number=page_number, is_first_page=is_first_page)
        
#         # Save generated content
#         with open(os.path.join(page_folder, "slide_script.txt"), "w", encoding="utf-8") as f:
#             f.write(slide_content)
        
#         with open(os.path.join(page_folder, "teacher_script.txt"), "w", encoding="utf-8") as f:
#             f.write(teacher_script)
        
#         # Extract images using exact image extraction code from reference
#         try:
#             # Convert PIL image to numpy array
#             image_np = np.array(page_image)
            
#             # Run layout detection
#             layout = lp_model.detect(image_np)
            
#             # Filter for figures and tables
#             blocks = [b for b in layout if b.type in ["Figure", "Table"]]
            
#             if blocks:
#                 # Extract coordinates and create rectangles
#                 rects = [tuple(map(int, b.coordinates)) for b in blocks]
                
#                 # Merge nearby rectangles
#                 merged_rects = merge_rectangles(rects, margin=20)
                
#                 # Save each detected region
#                 for j, rect in enumerate(merged_rects, start=1):
#                     x0, y0, x1, y1 = rect
#                     cropped = page_image.crop((x0, y0, x1, y1))
#                     cropped.save(os.path.join(page_folder, f"image_{j}.png"))
#                     print(f"    Saved image {j} from page {page_number}")
#             else:
#                 print(f"    No figure/table blocks detected on page {page_number}.")
#                 # Don't save anything if no figures/tables detected
#         except Exception as e:
#             print(f"Error in layout detection: {e}")
        
#         # Add small delay to avoid API rate limits
#         time.sleep(1)
        
#         print(f"Finished processing page {page_number}.")


# def main():
#     parser = argparse.ArgumentParser(description="Convert PDF to PowerPoint using GPT for content generation")
#     parser.add_argument("pdf_path", help="Path to the PDF file")
#     parser.add_argument("--output_dir", default="output", help="Output directory for processed files")
#     parser.add_argument("--pptx", default="output_presentation.pptx", help="Output PowerPoint filename")
#     parser.add_argument("--title_font_size", type=int, default=32, help="Font size for slide titles")
#     parser.add_argument("--content_font_size", type=int, default=14, help="Font size for slide content")
#     parser.add_argument("--widescreen", type=bool, default=True, help="Use 16:9 widescreen format (True) or 4:3 standard format (False)")
#     args = parser.parse_args()
    
#     if os.path.exists(args.output_dir):
#         print(f"Cleaning up existing output directory: {args.output_dir}")
#         shutil.rmtree(args.output_dir)
    
#     process_pdf(args.pdf_path, args.output_dir)
#     create_pptx(args.output_dir, pptx_filename=args.pptx, 
#                 title_font_size=args.title_font_size, 
#                 content_font_size=args.content_font_size)
#     print("Processing complete.")


# if __name__ == "__main__":
#     main()



import os
import argparse
import fitz
import numpy as np
from PIL import Image
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.dml.color import RGBColor
import openai
import time
import shutil
import warnings
import layoutparser as lp
from pptx.enum.text import PP_ALIGN

# Suppress warnings
warnings.filterwarnings("ignore")

# Handle PIL compatibility
if not hasattr(Image, "LINEAR"):
    Image.LINEAR = Image.BILINEAR

# Required pip packages:
# pip install python-pptx PyMuPDF pdf2image pillow openai layoutparser detectron2 torch

def extract_text_from_pdf(pdf_path):
    """Extract text from each page of the PDF"""
    doc = fitz.open(pdf_path)
    texts = []
    for page in doc:
        texts.append(page.get_text())
    return texts


def generate_with_gpt(text, prompt_type="slide", page_number=1, is_first_page=False, is_overview_page=False):
    """Generate content using GPT API"""
    # Configure your OpenAI API key
    openai.api_key = os.environ.get("OPENAI_API_KEY")
    
    if not openai.api_key:
        raise ValueError("Please set the OPENAI_API_KEY environment variable")
    
    # Define prompts for different content types
    prompts = {
        "slide": "Create 3-5 concise bullet points for a PowerPoint slide from the following text. "
                 "Each bullet point must be 10 words or less. Focus only on key concepts. "
                 "Format as bullet points with • symbol. "
                 "Also include a short, catchy title (up to 5 words) for this slide at the very beginning, "
                 "prefixed with 'TITLE: ':\n\n",
                 
        "discussion": "Based on the following content, generate ONE comprehensive discussion question "
                      "that encourages critical thinking and deeper exploration of the topic. "
                      "The question should be thought-provoking and open-ended. "
                      "Format as a single bullet point with • symbol:\n\n",
                 
        "speech": "Write a conversational lecture script for this slide as if you're an experienced professor speaking to your class. Build naturally on concepts covered earlier in your lecture and maintain a logical progression of ideas. Use a warm, natural teaching voice with occasional rhetorical questions, thoughtful pauses, and authentic phrases ('Now, what's interesting here is...', 'Let me emphasize this point...', 'You might be wondering...'). Include brief anecdotes or examples where appropriate, connect to previously explained ideas, and anticipate student questions. Vary your sentence structure with both short, punchy statements and more complex explanations. Limit to 150-200 words, maintain a clear focus on the slide content, and ensure it flows naturally when read aloud. The script should sound like authentic classroom teaching that fits seamlessly into the overall lecture:\n\n",

        "overview": "Create an overview slide that summarizes the main topics covered in this presentation. "
                    "Format as 3-5 concise bullet points with • symbol. "
                    "Each bullet point should be 10 words or less. "
    }
    
    # Special handling for first page
    if is_first_page and prompt_type == "slide":
        prompts["slide"] = "This is the FIRST SLIDE of the presentation. Create a title slide with: " + \
                           "1) A main title (5-7 words maximum) prefixed with 'MAIN_TITLE: ' " + \
                           "2) A subtitle (10-15 words) prefixed with 'SUBTITLE: ' " + \
                           "Both should capture the essence of the document. Based on this text:\n\n"
    
    # Limit input text length to avoid token limits
    max_input_length = 3000
    text = text[:max_input_length]
    
    try:
        # Send request to OpenAI API
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",  # You can also use "gpt-4" for better quality
            messages=[
                {"role": "system", "content": "You are a helpful assistant that creates PowerPoint content."},
                {"role": "user", "content": prompts[prompt_type] + text}
            ],
            max_tokens=1500,
            temperature=0.7
        )
        
        # Extract the generated content
        generated_content = response.choices[0].message.content.strip()
        return generated_content
    
    except Exception as e:
        print(f"Error generating content with GPT: {e}")
        # Fallback to basic extraction if GPT API fails
        if prompt_type == "slide":
            if is_first_page:
                return "MAIN_TITLE: Document Overview\nSUBTITLE: Key concepts and information"
            else:
                return "TITLE: Page " + str(page_number) + " Content\n" + extract_basic_bullet_points(text)
        else:
            return "Teacher notes for slide " + str(page_number) + ":\n\n" + text[:300] + "...\n\n[Note: Content truncated due to API error]"


def extract_basic_bullet_points(text, num_points=3):
    """Extract simple bullet points from text as a fallback when API fails"""
    sentences = []
    for sentence in text.split('.'):
        if len(sentence.strip()) > 10:  # Filter out very short fragments
            sentences.append(sentence.strip())
            if len(sentences) >= num_points * 3:  # Get more sentences than needed to select from
                break
    
    # Select a few sentences for bullet points
    selected = []
    step = max(1, len(sentences) // num_points)
    for i in range(0, min(len(sentences), num_points * step), step):
        if i < len(sentences):
            # Truncate long sentences
            point = sentences[i]
            if len(point) > 50:
                point = point[:47] + "..."
            selected.append("• " + point)
    
    # If we couldn't extract enough, add generic points
    while len(selected) < num_points:
        selected.append(f"• Key point {len(selected) + 1}")
    
    return "\n".join(selected)


def clean_title(title_text):
    """Remove asterisks from titles and clean up formatting"""
    return title_text.replace('*', '').strip()


# def create_pptx(output_dir, pptx_filename="output_presentation.pptx", title_font_size=32, content_font_size=14):
#     """Create PowerPoint presentation from processed content with adaptive layout"""
#     prs = Presentation()
    
#     # Change slide dimensions to 16:9 aspect ratio
#     prs.slide_width = int(Inches(16))
#     prs.slide_height = int(Inches(9))
    
#     blank_layout = prs.slide_layouts[6]  # Blank layout
#     title_layout = prs.slide_layouts[0]  # Title slide layout
    
#     # Calculate slide dimensions and margins - IMPROVED CENTERING
#     SLIDE_WIDTH = prs.slide_width / 914400
#     SLIDE_HEIGHT = prs.slide_height / 914400
#     left_margin = 1.0
#     right_margin = 1.0
#     top_margin = 1.5
#     bottom_margin = 0.8
#     title_spacing = 0.7
#     gutter_in = 0.5
#     usable_width = SLIDE_WIDTH - left_margin - right_margin
#     usable_height = SLIDE_HEIGHT - top_margin - bottom_margin
    
#     # Get and sort page folders
#     page_folders = []
#     for item in os.listdir(output_dir):
#         if item.startswith("page_") and os.path.isdir(os.path.join(output_dir, item)):
#             try:
#                 page_num = int(item.split("_")[1])
#                 page_folders.append((page_num, item))
#             except:
#                 continue
    
#     page_folders.sort()
    
#     # Process each page folder
#     for page_num, folder in page_folders:
#         folder_path = os.path.join(output_dir, folder)
#         is_first_page = (page_num == 1)
        
#         # Read generated slide content
#         slide_content_file = os.path.join(folder_path, "slide_script.txt")
#         slide_content = ""
#         slide_title = f"Page {page_num}"
        
#         if os.path.exists(slide_content_file):
#             try:
#                 with open(slide_content_file, "r", encoding="utf-8") as f:
#                     slide_content = f.read()
                    
#                     # Extract title from slide content
#                     if is_first_page and "MAIN_TITLE:" in slide_content:
#                         # Extract main title and subtitle for first slide
#                         main_title_match = slide_content.split("MAIN_TITLE:", 1)[1].split("\n", 1)[0].strip()
#                         main_title_match = clean_title(main_title_match)
                        
#                         subtitle_match = ""
#                         if "SUBTITLE:" in slide_content:
#                             subtitle_match = slide_content.split("SUBTITLE:", 1)[1].split("\n", 1)[0].strip()
#                             subtitle_match = clean_title(subtitle_match)
                        
#                         # Create title slide
#                         title_slide = prs.slides.add_slide(title_layout)
#                         title = title_slide.shapes.title
#                         subtitle = title_slide.placeholders[1]
                        
#                         # Add a background shape to maintain consistent visual style
#                         background = title_slide.shapes.add_shape(
#                             1,  # Rectangle
#                             0, 0, prs.slide_width, int(Inches(1.2))
#                         )
#                         background.fill.solid()
#                         background.fill.fore_color.rgb = RGBColor(0, 114, 198)
#                         background.line.fill.background()
                        
#                         # Bring title and subtitle to front
#                         title.element.getparent().remove(title.element)
#                         title_slide.shapes._spTree.append(title.element)
#                         subtitle.element.getparent().remove(subtitle.element)
#                         title_slide.shapes._spTree.append(subtitle.element)
                        
#                         # Positioning
#                         title.top = int(Inches(2.5))
#                         title.left = int(Inches(1.0))
#                         title.width = int(Inches(14.0))
#                         title.height = int(Inches(1.5))
#                         subtitle.top = int(Inches(4.0))
#                         subtitle.left = int(Inches(1.0))
#                         subtitle.width = int(Inches(14.0))
#                         subtitle.height = int(Inches(1.2))
                        
#                         # Text assignment and styling
#                         title.text = main_title_match
#                         title.text_frame.word_wrap = True
#                         title.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
#                         subtitle.text = subtitle_match
#                         subtitle.text_frame.word_wrap = True
#                         subtitle.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
#                         title.text_frame.paragraphs[0].runs[0].font.size = Pt(title_font_size + 8)
#                         subtitle.text_frame.paragraphs[0].runs[0].font.size = Pt(content_font_size + 4)
#                         title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
#                         subtitle.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
                        
#                         continue
#                     else:
#                         # Extract slide title for regular slides
#                         if "TITLE:" in slide_content:
#                             slide_title = slide_content.split("TITLE:", 1)[1].split("\n", 1)[0].strip()
#                             slide_title = clean_title(slide_title)
#                             slide_content = slide_content.split("\n", 1)[1] if "\n" in slide_content else ""
#             except:
#                 slide_content = f"Content for page {page_num}"
        
#         # Check for images
#         image_files = []
#         try:
#             image_files = sorted([
#                 f for f in os.listdir(folder_path) 
#                 if f.startswith("image_") and f.endswith(".png")
#             ])
#         except:
#             pass
        
#         # Create content slide
#         try:
#             slide = prs.slides.add_slide(blank_layout)
            
#             # Add title banner
#             title_banner = slide.shapes.add_shape(
#                 1,
#                 0, 0, prs.slide_width, int(Inches(1.0))
#             )
#             title_banner.fill.solid()
#             title_banner.fill.fore_color.rgb = RGBColor(0, 114, 198)
#             title_banner.line.fill.background()
            
#             # Title textbox
#             title_box = slide.shapes.add_textbox(
#                 Inches(left_margin),
#                 Inches(0.25),
#                 Inches(usable_width),
#                 Inches(0.7)
#             )
#             title_tf = title_box.text_frame
#             title_tf.word_wrap = True
#             title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
#             title_tf.text = slide_title
#             title_tf.paragraphs[0].alignment = 1
#             title_tf.paragraphs[0].runs[0].font.size = Pt(title_font_size)
#             title_tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
            
#             # If no images
#             if not image_files:
#                 content_box = slide.shapes.add_textbox(
#                     Inches(left_margin + 0.5),
#                     Inches(top_margin),
#                     Inches(usable_width - 1.0),
#                     Inches(usable_height)
#                 )
#                 content_tf = content_box.text_frame
#                 content_tf.text = slide_content
#                 content_tf.word_wrap = True
#                 content_tf.auto_size = MSO_AUTO_SIZE.NONE
#                 content_tf.vertical_anchor = MSO_ANCHOR.TOP
#                 larger_content_font_size = content_font_size + 6
#                 for paragraph in content_tf.paragraphs:
#                     for run in paragraph.runs:
#                         run.font.size = Pt(larger_content_font_size)
            
#             # If images available
#             elif image_files:
#                 img_path = os.path.join(folder_path, image_files[0])
#                 try:
#                     with Image.open(img_path) as pil_img:
#                         orig_w, orig_h = pil_img.size
#                     aspect_ratio = orig_h / orig_w
#                     if aspect_ratio < 0.75:
#                         content_height = usable_height * 0.40
#                         content_box = slide.shapes.add_textbox(
#                             Inches(left_margin + 0.5),
#                             Inches(top_margin),
#                             Inches(usable_width - 1.0),
#                             Inches(content_height)
#                         )
#                         pic_width_in = usable_width * 0.85
#                         pic_height_in = pic_width_in * aspect_ratio
#                         if pic_height_in > (usable_height - content_height - 0.3):
#                             pic_height_in = usable_height - content_height - 0.3
#                             pic_width_in = pic_height_in / aspect_ratio
#                         pic_left = left_margin + (usable_width - pic_width_in) / 2
#                         pic_top = top_margin + content_height + 0.3
#                     else:
#                         left_col_width = (usable_width - gutter_in) * 0.42
#                         right_col_width = (usable_width - gutter_in) * 0.58
#                         content_box = slide.shapes.add_textbox(
#                             Inches(left_margin + 0.3),
#                             Inches(top_margin),
#                             Inches(left_col_width - 0.3),
#                             Inches(usable_height)
#                         )
#                         pic_width_in = right_col_width * 0.95
#                         pic_height_in = pic_width_in * aspect_ratio
#                         if pic_height_in > usable_height:
#                             pic_height_in = usable_height * 0.95
#                             pic_width_in = pic_height_in / aspect_ratio
#                         pic_left = left_margin + left_col_width + gutter_in
#                         pic_top = top_margin + (usable_height - pic_height_in) / 2
#                     content_tf = content_box.text_frame
#                     content_tf.text = slide_content
#                     content_tf.word_wrap = True
#                     content_tf.auto_size = MSO_AUTO_SIZE.NONE
#                     content_tf.vertical_anchor = MSO_ANCHOR.TOP
#                     for paragraph in content_tf.paragraphs:
#                         for run in paragraph.runs:
#                             run.font.size = Pt(content_font_size + 2)
#                     slide.shapes.add_picture(
#                         img_path,
#                         Inches(pic_left),
#                         Inches(pic_top),
#                         width=Inches(pic_width_in),
#                         height=Inches(pic_height_in)
#                     )
#                 except Exception as e:
#                     print(f"Error processing image: {e}")
#                     content_box = slide.shapes.add_textbox(
#                         Inches(left_margin + 0.5),
#                         Inches(top_margin),
#                         Inches(usable_width - 1.0),
#                         Inches(usable_height)
#                     )
#                     content_tf = content_box.text_frame
#                     content_tf.text = slide_content
#                     content_tf.word_wrap = True
#                     for paragraph in content_tf.paragraphs:
#                         for run in paragraph.runs:
#                             run.font.size = Pt(content_font_size + 6)
#         except Exception as e:
#             print(f"Error creating slide for page {page_num}: {e}")
#             try:
#                 slide = prs.slides.add_slide(blank_layout)
#                 text_box = slide.shapes.add_textbox(
#                     Inches(left_margin + 0.5),
#                     Inches(top_margin),
#                     Inches(usable_width - 1.0),
#                     Inches(usable_height)
#                 )
#                 text_tf = text_box.text_frame
#                 text_tf.text = f"Page {page_num}\n\n" + slide_content
#                 for paragraph in text_tf.paragraphs:
#                     for run in paragraph.runs:
#                         run.font.size = Pt(content_font_size + 4)
#             except:
#                 print(f"Failed to create even basic slide for page {page_num}")
    
#     # -----------------新增讨论幻灯片-----------------
#     discussion_slide = prs.slides.add_slide(blank_layout)

#     # Title banner for discussion slide
#     title_banner = discussion_slide.shapes.add_shape(
#         1,
#         0, 0, prs.slide_width, int(Inches(1.0))
#     )
#     title_banner.fill.solid()
#     title_banner.fill.fore_color.rgb = RGBColor(0, 114, 198)
#     title_banner.line.fill.background()

#     # Discussion title textbox
#     title_box = discussion_slide.shapes.add_textbox(
#         Inches(left_margin),
#         Inches(0.25),
#         Inches(usable_width),
#         Inches(0.7)
#     )
#     title_tf = title_box.text_frame
#     title_tf.word_wrap = True
#     title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
#     title_tf.text = "Discussion"
#     title_tf.paragraphs[0].alignment = 1
#     title_tf.paragraphs[0].runs[0].font.size = Pt(title_font_size)
#     title_tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)

#     # Aggregate content from slides to generate context
#     aggregated_content = ""
#     for page_num, folder in page_folders:
#         slide_content_file = os.path.join(output_dir, folder, "slide_script.txt")
#         if os.path.exists(slide_content_file):
#             with open(slide_content_file, "r", encoding="utf-8") as f:
#                 aggregated_content += f.read() + "\n"

#     # Generate discussion questions dynamically using GPT
#     prompt = (
#         "Based on the following presentation content, generate 3 thoughtful discussion questions "
#         "that encourage critical thinking and deeper exploration of the topics. Format as bullet points with •:\n\n"
#         f"{aggregated_content[:3000]}"
#     )

#     discussion_questions_text = generate_with_gpt(prompt, prompt_type="slide")

#     # Remove title if present
#     discussion_questions_lines = discussion_questions_text.split('\n')
#     discussion_questions_lines = [line for line in discussion_questions_lines if line.strip().startswith('•')]
#     discussion_questions_text = '\n'.join(discussion_questions_lines)

#     content_box = discussion_slide.shapes.add_textbox(
#         Inches(left_margin + 0.5),
#         Inches(top_margin),
#         Inches(usable_width - 1.0),
#         Inches(usable_height)
#     )
#     content_tf = content_box.text_frame
#     content_tf.word_wrap = True
#     content_tf.auto_size = MSO_AUTO_SIZE.NONE
#     content_tf.vertical_anchor = MSO_ANCHOR.TOP
#     content_tf.text = discussion_questions_text

#     for paragraph in content_tf.paragraphs:
#         for run in paragraph.runs:
#             run.font.size = Pt(content_font_size + 2)
#     # ------------------------------------------------
    
#     # Save presentation
#     try:
#         output_directory = os.path.dirname(pptx_filename)
#         if output_directory and not os.path.exists(output_directory):
#             os.makedirs(output_directory)
#         prs.save(pptx_filename)
#         print(f"PowerPoint presentation saved as '{pptx_filename}'.")
#     except Exception as e:
#         print(f"Error saving presentation: {e}")
def create_pptx(output_dir, pptx_filename="output_presentation.pptx", title_font_size=32, content_font_size=14, author_name="", discussion_per_slide=False):
    """Create PowerPoint presentation from processed content with adaptive layout"""
    prs = Presentation()
    
    # Change slide dimensions to 16:9 aspect ratio
    prs.slide_width = int(Inches(16))
    prs.slide_height = int(Inches(9))
    
    blank_layout = prs.slide_layouts[6]  # Blank layout
    title_layout = prs.slide_layouts[0]  # Title slide layout
    
    # Calculate slide dimensions and margins - IMPROVED CENTERING
    SLIDE_WIDTH = prs.slide_width / 914400
    SLIDE_HEIGHT = prs.slide_height / 914400
    left_margin = 1.0
    right_margin = 1.0
    top_margin = 1.5
    bottom_margin = 0.8
    title_spacing = 0.7
    gutter_in = 0.5
    usable_width = SLIDE_WIDTH - left_margin - right_margin
    usable_height = SLIDE_HEIGHT - top_margin - bottom_margin
    
    # Get and sort page folders
    page_folders = []
    for item in os.listdir(output_dir):
        if item.startswith("page_") and os.path.isdir(os.path.join(output_dir, item)):
            try:
                page_num = int(item.split("_")[1])
                page_folders.append((page_num, item))
            except:
                continue
    
    page_folders.sort()
    
    # First, collect all texts to generate an overview
    all_text = ""
    for page_num, folder in page_folders:
        folder_path = os.path.join(output_dir, folder)
        text_file = os.path.join(folder_path, "page_text.txt")
        if os.path.exists(text_file):
            try:
                with open(text_file, "r", encoding="utf-8") as f:
                    all_text += f.read() + "\n\n"
            except:
                pass
    
    # Generate overview content
    overview_content = generate_with_gpt(all_text, prompt_type="overview", is_overview_page=True)
    
    # Process each page folder
    slide_index = 0
    
    # First, create title slide (first page)
    if page_folders and page_folders[0][0] == 1:
        page_num, folder = page_folders[0]
        folder_path = os.path.join(output_dir, folder)
        slide_content_file = os.path.join(folder_path, "slide_script.txt")
        slide_content = ""
        
        if os.path.exists(slide_content_file):
            try:
                with open(slide_content_file, "r", encoding="utf-8") as f:
                    slide_content = f.read()
                    
                # Extract main title and subtitle for first slide
                main_title_match = ""
                subtitle_match = ""
                
                if "MAIN_TITLE:" in slide_content:
                    main_title_match = slide_content.split("MAIN_TITLE:", 1)[1].split("\n", 1)[0].strip()
                    main_title_match = clean_title(main_title_match)
                
                if "SUBTITLE:" in slide_content:
                    subtitle_match = slide_content.split("SUBTITLE:", 1)[1].split("\n", 1)[0].strip()
                    subtitle_match = clean_title(subtitle_match)
                
                # Create title slide
                title_slide = prs.slides.add_slide(title_layout)
                title = title_slide.shapes.title
                subtitle = title_slide.placeholders[1]
                
                # Add a background shape to maintain consistent visual style
                background = title_slide.shapes.add_shape(
                    1,  # Rectangle
                    0, 0, prs.slide_width, int(Inches(1.2))
                )
                background.fill.solid()
                background.fill.fore_color.rgb = RGBColor(0, 114, 198)
                background.line.fill.background()
                
                # Bring title and subtitle to front
                title.element.getparent().remove(title.element)
                title_slide.shapes._spTree.append(title.element)
                subtitle.element.getparent().remove(subtitle.element)
                title_slide.shapes._spTree.append(subtitle.element)
                
                # Positioning
                title.top = int(Inches(2.5))
                title.left = int(Inches(1.0))
                title.width = int(Inches(14.0))
                title.height = int(Inches(1.5))
                subtitle.top = int(Inches(4.0))
                subtitle.left = int(Inches(1.0))
                subtitle.width = int(Inches(14.0))
                subtitle.height = int(Inches(1.2))
                
                # Text assignment and styling
                title.text = main_title_match
                title.text_frame.word_wrap = True
                title.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                subtitle.text = subtitle_match
                subtitle.text_frame.word_wrap = True
                subtitle.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                title.text_frame.paragraphs[0].runs[0].font.size = Pt(title_font_size + 8)
                subtitle.text_frame.paragraphs[0].runs[0].font.size = Pt(content_font_size + 4)
                title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
                subtitle.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
                
                # Add author name if provided
                if author_name:
                    author_box = title_slide.shapes.add_textbox(
                        int(Inches(1.0)),
                        int(Inches(5.2)),
                        int(Inches(14.0)),
                        int(Inches(0.5))
                    )
                    author_tf = author_box.text_frame
                    author_tf.text = f"By: {author_name}"
                    author_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

                    if not author_tf.paragraphs[0].runs:
                        p = author_tf.paragraphs[0]
                        p.add_run()
                        
                    author_tf.paragraphs[0].runs[0].font.size = Pt(content_font_size)
                    author_tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
                
                slide_index += 1
            except Exception as e:
                print(f"Error creating title slide: {e}")
    
    # Now add the overview slide
    try:
        overview_slide = prs.slides.add_slide(blank_layout)
        
        # Add title banner
        title_banner = overview_slide.shapes.add_shape(
            1,
            0, 0, prs.slide_width, int(Inches(1.0))
        )
        title_banner.fill.solid()
        title_banner.fill.fore_color.rgb = RGBColor(0, 114, 198)
        title_banner.line.fill.background()
        
        # Title textbox
        title_box = overview_slide.shapes.add_textbox(
            Inches(left_margin),
            Inches(0.25),
            Inches(usable_width),
            Inches(0.7)
        )
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        title_tf.text = "Presentation Overview"
        title_tf.paragraphs[0].alignment = 1
        title_tf.paragraphs[0].runs[0].font.size = Pt(title_font_size)
        title_tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Content box
        content_box = overview_slide.shapes.add_textbox(
            Inches(left_margin + 0.5),
            Inches(top_margin),
            Inches(usable_width - 1.0),
            Inches(usable_height)
        )
        content_tf = content_box.text_frame
        
        # Clean the overview content to remove any "TITLE:" prefix
        cleaned_content = overview_content
        if "TITLE:" in overview_content:
            # Remove the entire line containing "TITLE:"
            cleaned_content = "\n".join([line for line in overview_content.split('\n') 
                                        if "TITLE:" not in line])
        
        # Add spacing between bullet points in overview
        if "•" in cleaned_content:
            overview_lines = cleaned_content.split('\n')
            processed_lines = []
            for line in overview_lines:
                if line.strip():  # Skip empty lines
                    processed_lines.append(line)
                    if line.strip().startswith('•'):
                        processed_lines.append('')  # Add an empty line after each bullet point
            processed_content = '\n'.join(processed_lines)
            content_tf.text = processed_content
        else:
            content_tf.text = cleaned_content
            
        content_tf.word_wrap = True
        content_tf.auto_size = MSO_AUTO_SIZE.NONE
        content_tf.vertical_anchor = MSO_ANCHOR.TOP
        
        # Format all paragraphs in the content
        for paragraph in content_tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(content_font_size + 6)
        
        slide_index += 1
    except Exception as e:
        print(f"Error creating overview slide: {e}")
    
    # Process content slides (skip the first page as it's already handled)
    for page_num, folder in page_folders:
        if page_num == 1:  # Skip first page as it's already processed as title slide
            continue
            
        folder_path = os.path.join(output_dir, folder)
        
        # Read generated slide content
        slide_content_file = os.path.join(folder_path, "slide_script.txt")
        slide_content = ""
        slide_title = f"Page {page_num}"
        original_content = ""
        
        if os.path.exists(slide_content_file):
            try:
                with open(slide_content_file, "r", encoding="utf-8") as f:
                    slide_content = f.read()
                    original_content = slide_content  # Store original content for discussion question generation
                    
                    # Extract slide title for regular slides
                    if "TITLE:" in slide_content:
                        slide_title = slide_content.split("TITLE:", 1)[1].split("\n", 1)[0].strip()
                        slide_title = clean_title(slide_title)
                        slide_content = slide_content.split("\n", 1)[1] if "\n" in slide_content else ""
            except:
                slide_content = f"Content for page {page_num}"
                original_content = slide_content
        
        # Check for images
        image_files = []
        try:
            image_files = sorted([
                f for f in os.listdir(folder_path) 
                if f.startswith("image_") and f.endswith(".png")
            ])
        except:
            pass
        
        # Create content slide
        try:
            slide = prs.slides.add_slide(blank_layout)
            
            # Add title banner
            title_banner = slide.shapes.add_shape(
                1,
                0, 0, prs.slide_width, int(Inches(1.0))
            )
            title_banner.fill.solid()
            title_banner.fill.fore_color.rgb = RGBColor(0, 114, 198)
            title_banner.line.fill.background()
            
            # Title textbox
            title_box = slide.shapes.add_textbox(
                Inches(left_margin),
                Inches(0.25),
                Inches(usable_width),
                Inches(0.7)
            )
            title_tf = title_box.text_frame
            title_tf.word_wrap = True
            title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            title_tf.text = slide_title
            title_tf.paragraphs[0].alignment = 1
            title_tf.paragraphs[0].runs[0].font.size = Pt(title_font_size)
            title_tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
            
            # If no images
            if not image_files:
                content_box = slide.shapes.add_textbox(
                    Inches(left_margin + 0.5),
                    Inches(top_margin),
                    Inches(usable_width - 1.0),
                    Inches(usable_height)
                )
                content_tf = content_box.text_frame
                
                if "•" in slide_content:
                    content_lines = slide_content.split('\n')
                    processed_lines = []
                    for line in content_lines:
                        processed_lines.append(line)
                        if line.strip().startswith('•'):
                            processed_lines.append('')  # Add an empty line after each bullet point
                    slide_content = '\n'.join(processed_lines)
                    
                content_tf.text = slide_content
                content_tf.word_wrap = True
                content_tf.auto_size = MSO_AUTO_SIZE.NONE
                content_tf.vertical_anchor = MSO_ANCHOR.TOP
                larger_content_font_size = content_font_size + 6
                for paragraph in content_tf.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(larger_content_font_size)
            
            # If images available
            elif image_files:
                img_path = os.path.join(folder_path, image_files[0])
                try:
                    with Image.open(img_path) as pil_img:
                        orig_w, orig_h = pil_img.size
                    aspect_ratio = orig_h / orig_w
                        
                    if aspect_ratio < 0.75:
                        # Wider image - text on top, image below
                        content_height = usable_height * 0.40
                            
                        content_box = slide.shapes.add_textbox(
                            Inches(left_margin + 0.5),
                            Inches(top_margin),
                            Inches(usable_width - 1.0),
                            Inches(content_height)
                        )
                        pic_width_in = usable_width * 0.85
                        pic_height_in = pic_width_in * aspect_ratio
                        if pic_height_in > (usable_height - content_height - 0.3):
                            pic_height_in = usable_height - content_height - 0.3
                            pic_width_in = pic_height_in / aspect_ratio
                        pic_left = left_margin + (usable_width - pic_width_in) / 2
                        pic_top = top_margin + content_height + 0.3
                    else:
                        # Taller image - text on left, image on right
                        left_col_width = (usable_width - gutter_in) * 0.42
                            
                        right_col_width = (usable_width - gutter_in) - left_col_width
                        content_box = slide.shapes.add_textbox(
                            Inches(left_margin + 0.3),
                            Inches(top_margin),
                            Inches(left_col_width - 0.3),
                            Inches(usable_height)
                        )
                        pic_width_in = right_col_width * 0.95
                        pic_height_in = pic_width_in * aspect_ratio
                        if pic_height_in > usable_height:
                            pic_height_in = usable_height * 0.95
                            pic_width_in = pic_height_in / aspect_ratio
                        pic_left = left_margin + left_col_width + gutter_in
                        pic_top = top_margin + (usable_height - pic_height_in) / 2
                    
                    content_tf = content_box.text_frame
                    content_tf.text = slide_content
                    content_tf.word_wrap = True
                    content_tf.auto_size = MSO_AUTO_SIZE.NONE
                    content_tf.vertical_anchor = MSO_ANCHOR.TOP
                    for paragraph in content_tf.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(content_font_size + 2)
                    slide.shapes.add_picture(
                        img_path,
                        Inches(pic_left),
                        Inches(pic_top),
                        width=Inches(pic_width_in),
                        height=Inches(pic_height_in)
                    )
                except Exception as e:
                    print(f"Error processing image: {e}")
                    content_box = slide.shapes.add_textbox(
                        Inches(left_margin + 0.5),
                        Inches(top_margin),
                        Inches(usable_width - 1.0),
                        Inches(usable_height)
                    )
                    content_tf = content_box.text_frame
                    content_tf.text = slide_content
                    content_tf.word_wrap = True
                    for paragraph in content_tf.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(content_font_size + 6)
            
            slide_index += 1
            
            # Generate discussion question and create a separate discussion slide if requested
            if discussion_per_slide:
                # Use the original text from the page to generate a discussion question
                page_text_file = os.path.join(folder_path, "page_text.txt")
                page_text = ""
                if os.path.exists(page_text_file):
                    try:
                        with open(page_text_file, "r", encoding="utf-8") as f:
                            page_text = f.read()
                    except:
                        page_text = original_content
                        
                discussion_prompt = (
                    "Based on the following content, generate ONE comprehensive discussion question "
                    "that encourages critical thinking and deeper exploration of the topic. "
                    "The question should be thought-provoking and open-ended. "
                    "Format as a single bullet point with • symbol:\n\n"
                    f"{page_text[:1500]}"
                )
                
                discussion_content = generate_with_gpt(discussion_prompt, prompt_type="discussion")
                
                # Create the discussion slide
                discussion_slide = prs.slides.add_slide(blank_layout)
                
                # Add title banner for discussion slide
                disc_title_banner = discussion_slide.shapes.add_shape(
                    1,
                    0, 0, prs.slide_width, int(Inches(1.0))
                )
                disc_title_banner.fill.solid()
                disc_title_banner.fill.fore_color.rgb = RGBColor(0, 114, 198)
                disc_title_banner.line.fill.background()
                
                # Title textbox for discussion slide
                disc_title_box = discussion_slide.shapes.add_textbox(
                    Inches(left_margin),
                    Inches(0.25),
                    Inches(usable_width),
                    Inches(0.7)
                )
                disc_title_tf = disc_title_box.text_frame
                disc_title_tf.word_wrap = True
                disc_title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                disc_title_tf.text = f"Discussion: {slide_title}"
                disc_title_tf.paragraphs[0].alignment = 1
                disc_title_tf.paragraphs[0].runs[0].font.size = Pt(title_font_size)
                disc_title_tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
                
                # Content box for discussion question
                disc_content_box = discussion_slide.shapes.add_textbox(
                    Inches(left_margin + 0.5),
                    Inches(top_margin),
                    Inches(usable_width - 1.0),
                    Inches(usable_height)
                )
                disc_content_tf = disc_content_box.text_frame
                disc_content_tf.word_wrap = True
                disc_content_tf.auto_size = MSO_AUTO_SIZE.NONE
                disc_content_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                # Process and format the discussion question
                if "•" in discussion_content:
                    disc_lines = [line.strip() for line in discussion_content.split('\n') 
                                if line.strip().startswith('•')]
                    if disc_lines:
                        discussion_text = "\n\n".join(disc_lines)
                        disc_content_tf.text = discussion_text
                    else:
                        disc_content_tf.text = "• What insights can we draw from this topic that might impact our understanding or approach?"
                else:
                    # If no bullet points found, use the raw content or a fallback
                    if '?' in discussion_content:
                        # Find the first question in the content
                        for line in discussion_content.split('\n'):
                            if '?' in line:
                                disc_content_tf.text = f"• {line.strip()}"
                                break
                    else:
                        disc_content_tf.text = f"• {discussion_content.strip()}"
                
                # Format text
                for paragraph in disc_content_tf.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER  # Center-align the question
                    for run in paragraph.runs:
                        run.font.size = Pt(content_font_size + 8)  # Make it larger than normal content
                        run.font.bold = True  # Bold for emphasis
                
                slide_index += 1
                
        except Exception as e:
            print(f"Error creating slide for page {page_num}: {e}")
            try:
                slide = prs.slides.add_slide(blank_layout)
                text_box = slide.shapes.add_textbox(
                    Inches(left_margin + 0.5),
                    Inches(top_margin),
                    Inches(usable_width - 1.0),
                    Inches(usable_height)
                )
                text_tf = text_box.text_frame
                text_tf.text = f"Page {page_num}\n\n" + slide_content
                for paragraph in text_tf.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(content_font_size + 4)
                
                slide_index += 1
            except:
                print(f"Failed to create even basic slide for page {page_num}")
    
    # Save presentation
    try:
        output_directory = os.path.dirname(pptx_filename)
        if output_directory and not os.path.exists(output_directory):
            os.makedirs(output_directory)
        prs.save(pptx_filename)
        print(f"PowerPoint presentation saved as '{pptx_filename}'.")
    except Exception as e:
        print(f"Error saving presentation: {e}")

def boxes_are_close(box1, box2, margin=10):
    x0, y0, x1, y1 = box1
    a0, b0, a1, b1 = box2
    gap_x = 0
    if x1 < a0:
        gap_x = a0 - x1
    elif a1 < x0:
        gap_x = x0 - a1
    gap_y = 0
    if y1 < b0:
        gap_y = b0 - y1
    elif b1 < y0:
        gap_y = y0 - b1
    return (gap_x <= margin) and (gap_y <= margin)


def union_box(box1, box2):
    x0, y0, x1, y1 = box1
    a0, b0, a1, b1 = box2
    return (min(x0, a0), min(y0, b0), max(x1, a1), max(y1, b1))


def merge_rectangles(rectangles, margin=10):
    if not rectangles:
        return []
    merged = rectangles[:]
    changed = True
    while changed:
        changed = False
        new_merged = []
        used = [False] * len(merged)
        for i in range(len(merged)):
            if used[i]:
                continue
            current = merged[i]
            for j in range(i + 1, len(merged)):
                if used[j]:
                    continue
                if boxes_are_close(current, merged[j], margin):
                    current = union_box(current, merged[j])
                    used[j] = True
                    changed = True
            new_merged.append(current)
        merged = new_merged
    return merged


def process_pdf(pdf_path, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    print("Converting PDF pages to images...")
    pages = convert_from_path(pdf_path, dpi=300)
    print(f"Converted {len(pages)} pages.")
    print("Extracting text from PDF...")
    texts = extract_text_from_pdf(pdf_path)
    from layoutparser.models import Detectron2LayoutModel
    print("Loading layout detection model...")
    lp_model = Detectron2LayoutModel(
        "lp://PubLayNet/faster_rcnn_R_50_FPN_3x/config",
        model_path="/models/layoutparser/publaynet/model_final_new.pth",
        extra_config=["MODEL.ROI_HEADS.SCORE_THRESH_TEST", 0.5],
        label_map={0: "Text", 1: "Title", 2: "List", 3: "Table", 4: "Figure"}
    )
    for i, page_image in enumerate(pages):
        page_number = i + 1
        is_first_page = (page_number == 1)
        print(f"\nProcessing page {page_number}...")
        page_folder = os.path.join(output_dir, f"page_{page_number}")
        os.makedirs(page_folder, exist_ok=True)
        text = texts[i] if i < len(texts) else ""
        with open(os.path.join(page_folder, "page_text.txt"), "w", encoding="utf-8") as f:
            f.write(text)
        print(f"  Generating slide script for page {page_number}...")
        slide_content = generate_with_gpt(text, prompt_type="slide", page_number=page_number, is_first_page=is_first_page)
        print(f"  Generating teacher script for page {page_number}...")
        teacher_script = generate_with_gpt(text, prompt_type="speech", page_number=page_number, is_first_page=is_first_page)
        with open(os.path.join(page_folder, "slide_script.txt"), "w", encoding="utf-8") as f:
            f.write(slide_content)
        with open(os.path.join(page_folder, "teacher_script.txt"), "w", encoding="utf-8") as f:
            f.write(teacher_script)
        try:
            image_np = np.array(page_image)
            layout = lp_model.detect(image_np)
            blocks = [b for b in layout if b.type in ["Figure", "Table"]]
            if blocks:
                rects = [tuple(map(int, b.coordinates)) for b in blocks]
                merged_rects = merge_rectangles(rects, margin=20)
                for j, rect in enumerate(merged_rects, start=1):
                    x0, y0, x1, y1 = rect
                    cropped = page_image.crop((x0, y0, x1, y1))
                    cropped.save(os.path.join(page_folder, f"image_{j}.png"))
                    print(f"    Saved image {j} from page {page_number}")
            else:
                print(f"    No figure/table blocks detected on page {page_number}.")
        except Exception as e:
            print(f"Error in layout detection: {e}")
        image_files = [
            f for f in os.listdir(page_folder) 
            if f.startswith("image_") and f.endswith(".png")
        ]
        
        if not image_files:
            print(f"  No images found for page {page_number}. Generating AI image...")
            prompt_text = text.strip()[:500]  # 限制长度
            if not prompt_text:
                prompt_text = "Create an abstract background for a professional presentation slide."
        time.sleep(1)
        print(f"Finished processing page {page_number}.")

# from ultralytics import YOLO 
# import cv2

# def process_pdf(pdf_path, output_dir):
#     os.makedirs(output_dir, exist_ok=True)
#     print("Converting PDF pages to images...")
#     pages = convert_from_path(pdf_path, dpi=300)
#     print(f"Converted {len(pages)} pages.")
#     print("Extracting text from PDF...")
#     texts = extract_text_from_pdf(pdf_path)
#     yolo_model = YOLO("model/yolo11n.pt")
#     for i, page_image in enumerate(pages):
#         page_number = i + 1
#         is_first_page = (page_number == 1)
#         print(f"\nProcessing page {page_number}...")
#         page_folder = os.path.join(output_dir, f"page_{page_number}")
#         os.makedirs(page_folder, exist_ok=True)
#         text = texts[i] if i < len(texts) else ""
#         with open(os.path.join(page_folder, "page_text.txt"), "w", encoding="utf-8") as f:
#             f.write(text)
#         print(f"  Generating slide script for page {page_number}...")
#         slide_content = generate_with_gpt(text, prompt_type="slide", page_number=page_number, is_first_page=is_first_page)
#         print(f"  Generating teacher script for page {page_number}...")
#         teacher_script = generate_with_gpt(text, prompt_type="speech", page_number=page_number, is_first_page=is_first_page)
#         with open(os.path.join(page_folder, "slide_script.txt"), "w", encoding="utf-8") as f:
#             f.write(slide_content)
#         with open(os.path.join(page_folder, "teacher_script.txt"), "w", encoding="utf-8") as f:
#             f.write(teacher_script)
#         try:
#             image_np = np.array(page_image)
#             image_bgr = cv2.cvtColor(image_np, cv2.COLOR_RGB2BGR)
#             results = yolo_model.predict(source=image_bgr, save = False, verbose=False)[0]
#             blocks = []
#             for box in results.boxes:
#                 cls_id = int(box.cls[0])
#                 label = yolo_model.names[cls_id]
#                 if label in ["Figure", "Table"]:
#                     x0, y0, x1, y1 = map(int, box.xyxy[0])
#                     blocks.append((x0, y0, x1, y1))
#             if blocks:
#                 merged_rects = merge_rectangles(blocks, margin=20)
#                 for j, rect in enumerate(merged_rects, start=1):
#                     x0, y0, x1, y1 = rect
#                     cropped = page_image.crop((x0, y0, x1, y1))
#                     cropped.save(os.path.join(page_folder, f"image_{j}.png"))
#                     print(f"    Saved image {j} from page {page_number}")
#             else:
#                 print(f"    No figure/table blocks detected on page {page_number}.")
#         except Exception as e:
#             print(f"Error in layout detection: {e}")
#         time.sleep(1)
#         print(f"Finished processing page {page_number}.")

def main():
    parser = argparse.ArgumentParser(description="Convert PDF to PowerPoint using GPT for content generation")
    parser.add_argument("pdf_path", help="Path to the PDF file")
    parser.add_argument("--output_dir", default="output", help="Output directory for processed files")
    parser.add_argument("--pptx", default="output_presentation.pptx", help="Output PowerPoint filename")
    parser.add_argument("--title_font_size", type=int, default=32, help="Font size for slide titles")
    parser.add_argument("--content_font_size", type=int, default=20, help="Font size for slide content")
    parser.add_argument("--widescreen", type=bool, default=True, help="Use 16:9 widescreen format (True) or 4:3 standard format (False)")
    parser.add_argument("--author", default="MAIL AI LAB MEMBERS", help="Author name to display on the title slide")
    parser.add_argument("--discussion_per_slide", action="store_true", default=True,
                    help="Add discussion questions to each content slide")
    args = parser.parse_args()
    if os.path.exists(args.output_dir):
        print(f"Cleaning up existing output directory: {args.output_dir}")
        shutil.rmtree(args.output_dir)
    process_pdf(args.pdf_path, args.output_dir)
    create_pptx(args.output_dir, pptx_filename=args.pptx, 
                title_font_size=args.title_font_size, 
                content_font_size=args.content_font_size,
                author_name=args.author,
                discussion_per_slide=args.discussion_per_slide)
    print("Processing complete.")


if __name__ == "__main__":
    main()